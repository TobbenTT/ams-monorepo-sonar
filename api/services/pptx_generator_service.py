"""
PPTX Generator Service
======================
Generates professional PowerPoint presentations from executive report data
using python-pptx. Outputs a styled .pptx file with KPIs, tables, budget
variance, critical equipment, and executive summary slides.
"""

import logging
import os
import tempfile
from datetime import datetime

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

log = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Style constants
# ---------------------------------------------------------------------------
DARK_BLUE = RGBColor(0x0F, 0x17, 0x2A)
PRIMARY_GREEN = RGBColor(0x04, 0x78, 0x57)
RED = RGBColor(0xDC, 0x26, 0x26)
ORANGE = RGBColor(0xEA, 0x58, 0x0C)
YELLOW = RGBColor(0xCA, 0x8A, 0x04)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF1, 0xF5, 0xF9)


# ---------------------------------------------------------------------------
# Internal helpers
# ---------------------------------------------------------------------------

def _safe_get(data: dict | None, key: str, default=None):
    """Safely retrieve a value from a dict, returning *default* if missing."""
    if data is None:
        return default
    return data.get(key, default)


def _fmt_number(value, decimals: int = 0) -> str:
    """Format a numeric value as a string, defaulting to 'N/A'."""
    if value is None:
        return "N/A"
    try:
        if decimals == 0:
            return f"{int(value):,}"
        return f"{float(value):,.{decimals}f}"
    except (ValueError, TypeError):
        return "N/A"


def _fmt_pct(value) -> str:
    """Format a percentage value."""
    if value is None:
        return "N/A"
    try:
        return f"{float(value):.1f}%"
    except (ValueError, TypeError):
        return "N/A"


def _fmt_currency(value) -> str:
    """Format a currency value."""
    if value is None:
        return "N/A"
    try:
        return f"${float(value):,.0f}"
    except (ValueError, TypeError):
        return "N/A"


def _set_cell_text(cell, text: str, font_size: int = 10, bold: bool = False,
                   alignment=PP_ALIGN.CENTER, font_color: RGBColor | None = None):
    """Set text and formatting on a table cell."""
    cell.text = str(text)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.bold = bold
            if font_color is not None:
                run.font.color.rgb = font_color


def _set_cell_bg(cell, color: RGBColor):
    """Set background fill on a table cell."""
    fill = cell.fill
    fill.solid()
    fill.fore_color.rgb = color


def _style_header_row(table, col_count: int, font_size: int = 10):
    """Style the first row of a table as the header (green bg, white text)."""
    for col_idx in range(col_count):
        cell = table.cell(0, col_idx)
        _set_cell_bg(cell, PRIMARY_GREEN)
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = WHITE
                run.font.bold = True
                run.font.size = Pt(font_size)


def _apply_alternating_rows(table, row_count: int, col_count: int,
                            start_row: int = 1):
    """Apply alternating white / light-gray background to data rows."""
    for row_idx in range(start_row, row_count):
        bg = LIGHT_GRAY if row_idx % 2 == 0 else WHITE
        for col_idx in range(col_count):
            _set_cell_bg(table.cell(row_idx, col_idx), bg)


def _variance_color(value) -> RGBColor:
    """Return traffic-light colour for a variance percentage."""
    try:
        v = abs(float(value))
    except (ValueError, TypeError):
        return WHITE
    if v > 10:
        return RED
    if v > 5:
        return YELLOW
    return GREEN


def _reactive_color(value) -> RGBColor:
    """Return traffic-light colour for reactive ratio percentage."""
    try:
        v = float(value)
    except (ValueError, TypeError):
        return WHITE
    if v > 30:
        return RED
    if v > 20:
        return YELLOW
    return GREEN


def _criticality_color(criticality: str) -> RGBColor | None:
    """Return background colour for criticality level."""
    c = (criticality or "").strip().upper()
    if c == "AA":
        return RED
    if c == "A":
        return ORANGE
    return None


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def _add_title_slide(prs: Presentation, report_data: dict):
    """Slide 1: Title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    plant_id = _safe_get(report_data, "plant_id", "N/A")
    period = _safe_get(report_data, "period", "N/A")
    generated_at = _safe_get(report_data, "generated_at", "")
    try:
        date_label = datetime.fromisoformat(generated_at).strftime("%d/%m/%Y")
    except (ValueError, TypeError):
        date_label = datetime.now().strftime("%d/%m/%Y")

    slide.shapes.title.text = "Reporte Ejecutivo de Mantenimiento"
    if slide.placeholders[1]:
        slide.placeholders[1].text = f"{plant_id} \u2014 {period} \u2014 {date_label}"


def _add_kpi_slide(prs: Presentation, report_data: dict):
    """Slide 2: KPI cards as a 3x4 table."""
    slides_data = _safe_get(report_data, "slides_data", [])
    kpi_data: dict = {}
    for sd in slides_data:
        if _safe_get(sd, "type") == "kpi_cards":
            kpi_data = _safe_get(sd, "data", {})
            break
    if not kpi_data:
        kpi_data = _safe_get(report_data, "kpis", {})

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "KPIs de Mantenimiento"

    rows, cols = 3, 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.8),
                                         Inches(9.0), Inches(2.5))
    table = table_shape.table

    # Row 0 — headers
    headers = ["OTs Totales", "OTs Completadas", "Correctivas", "Preventivas"]
    for i, h in enumerate(headers):
        _set_cell_text(table.cell(0, i), h, font_size=11, bold=True)
    _style_header_row(table, cols, font_size=11)

    # Row 1 — values
    vals = [
        _fmt_number(_safe_get(kpi_data, "total_work_orders", 0)),
        _fmt_number(_safe_get(kpi_data, "completed_work_orders", 0)),
        _fmt_number(_safe_get(kpi_data, "corrective_count", 0)),
        _fmt_number(_safe_get(kpi_data, "preventive_count", 0)),
    ]
    for i, v in enumerate(vals):
        _set_cell_text(table.cell(1, i), v, font_size=14, bold=True)

    # Row 2 — reactive ratio + avg hours
    reactive_pct = _safe_get(kpi_data, "reactive_ratio_pct", 0)
    avg_hours = _safe_get(kpi_data, "avg_completion_hours", 0)
    _set_cell_text(table.cell(2, 0), "Ratio Reactivo", font_size=10, bold=True)
    _set_cell_text(table.cell(2, 1), _fmt_pct(reactive_pct), font_size=14, bold=True)
    _set_cell_text(table.cell(2, 2), "Hrs Promedio", font_size=10, bold=True)
    _set_cell_text(table.cell(2, 3), _fmt_number(avg_hours, decimals=1), font_size=14, bold=True)

    # Traffic-light on reactive ratio cell
    color = _reactive_color(reactive_pct)
    _set_cell_bg(table.cell(2, 1), color)
    # Ensure text is visible on colored background
    for paragraph in table.cell(2, 1).text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = WHITE

    _apply_alternating_rows(table, rows, cols, start_row=1)
    # Re-apply the reactive cell colour (alternating rows may have overridden it)
    _set_cell_bg(table.cell(2, 1), color)
    for paragraph in table.cell(2, 1).text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = WHITE


def _add_backlog_slide(prs: Presentation, report_data: dict):
    """Slide 3: Backlog status table."""
    slides_data = _safe_get(report_data, "slides_data", [])
    backlog_data: dict = {}
    for sd in slides_data:
        if _safe_get(sd, "type") == "table" and "backlog" in _safe_get(sd, "title", "").lower():
            backlog_data = _safe_get(sd, "data", {})
            break
    if not backlog_data:
        backlog_data = _safe_get(report_data, "backlog_stats", {})

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Estado del Backlog"

    by_priority = _safe_get(backlog_data, "by_priority", {})
    total = _safe_get(backlog_data, "total", 0) or 0
    avg_age = _safe_get(backlog_data, "avg_age_days", 0)
    priorities = ["P1", "P2", "P3", "P4"]

    num_rows = 2 + len(priorities)  # header + priorities + total
    cols = 3
    table_shape = slide.shapes.add_table(num_rows, cols, Inches(1.5), Inches(1.8),
                                         Inches(7.0), Inches(0.4 * num_rows + 0.5))
    table = table_shape.table

    # Header
    for i, h in enumerate(["Prioridad", "Cantidad", "% del Total"]):
        _set_cell_text(table.cell(0, i), h, font_size=11, bold=True)
    _style_header_row(table, cols, font_size=11)

    # Priority rows
    for row_idx, p in enumerate(priorities, start=1):
        count = _safe_get(by_priority, p, 0) or 0
        pct = (count / total * 100) if total > 0 else 0
        _set_cell_text(table.cell(row_idx, 0), p, font_size=10, bold=True)
        _set_cell_text(table.cell(row_idx, 1), _fmt_number(count), font_size=10)
        _set_cell_text(table.cell(row_idx, 2), _fmt_pct(pct), font_size=10)

    # Total row
    last = len(priorities) + 1
    _set_cell_text(table.cell(last, 0), "Total", font_size=10, bold=True)
    _set_cell_text(table.cell(last, 1), _fmt_number(total), font_size=10, bold=True)
    _set_cell_text(table.cell(last, 2), f"Edad Promedio: {_fmt_number(avg_age, decimals=1)} dias",
                   font_size=10)

    _apply_alternating_rows(table, num_rows, cols)


def _add_budget_slide(prs: Presentation, report_data: dict):
    """Slide 4: Budget vs Actual variance table."""
    slides_data = _safe_get(report_data, "slides_data", [])
    budget_data: dict = {}
    for sd in slides_data:
        if _safe_get(sd, "type") == "chart":
            budget_data = _safe_get(sd, "data", {})
            break
    if not budget_data:
        budget_data = _safe_get(report_data, "budget_variance", {})

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Presupuesto vs Real"

    by_type = _safe_get(budget_data, "by_type", {})
    total_budget = _safe_get(budget_data, "total_budget", 0)
    total_actual = _safe_get(budget_data, "total_actual", 0)
    total_variance = _safe_get(budget_data, "variance_pct", 0)

    type_keys = list(by_type.keys())
    num_rows = 2 + len(type_keys)  # header + types + total
    cols = 4
    table_shape = slide.shapes.add_table(num_rows, cols, Inches(1.0), Inches(1.8),
                                         Inches(8.0), Inches(0.4 * num_rows + 0.5))
    table = table_shape.table

    # Header
    for i, h in enumerate(["Tipo OT", "Presupuesto", "Real", "Variacion %"]):
        _set_cell_text(table.cell(0, i), h, font_size=11, bold=True)
    _style_header_row(table, cols, font_size=11)

    # Type rows
    for row_idx, wo_type in enumerate(type_keys, start=1):
        entry = by_type[wo_type]
        if isinstance(entry, dict):
            bgt = _safe_get(entry, "budget", 0)
            act = _safe_get(entry, "actual", 0)
            var = _safe_get(entry, "variance_pct", 0)
        else:
            bgt = entry
            act = 0
            var = 0
        _set_cell_text(table.cell(row_idx, 0), str(wo_type), font_size=10, bold=True,
                       alignment=PP_ALIGN.LEFT)
        _set_cell_text(table.cell(row_idx, 1), _fmt_currency(bgt), font_size=10)
        _set_cell_text(table.cell(row_idx, 2), _fmt_currency(act), font_size=10)
        _set_cell_text(table.cell(row_idx, 3), _fmt_pct(var), font_size=10)
        # Colour variance cell
        _set_cell_bg(table.cell(row_idx, 3), _variance_color(var))
        for paragraph in table.cell(row_idx, 3).text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = WHITE

    # Total row
    last = len(type_keys) + 1
    _set_cell_text(table.cell(last, 0), "Total", font_size=10, bold=True,
                   alignment=PP_ALIGN.LEFT)
    _set_cell_text(table.cell(last, 1), _fmt_currency(total_budget), font_size=10, bold=True)
    _set_cell_text(table.cell(last, 2), _fmt_currency(total_actual), font_size=10, bold=True)
    _set_cell_text(table.cell(last, 3), _fmt_pct(total_variance), font_size=10, bold=True)
    _set_cell_bg(table.cell(last, 3), _variance_color(total_variance))
    for paragraph in table.cell(last, 3).text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.color.rgb = WHITE

    _apply_alternating_rows(table, num_rows, cols)
    # Re-apply variance colours (alternating rows may have overridden them)
    for row_idx in range(1, num_rows):
        cell = table.cell(row_idx, 3)
        text = cell.text_frame.paragraphs[0].runs[0].text if cell.text_frame.paragraphs and cell.text_frame.paragraphs[0].runs else "0"
        try:
            v = float(text.replace("%", "").replace(",", ""))
        except (ValueError, TypeError):
            v = 0
        _set_cell_bg(cell, _variance_color(v))
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.color.rgb = WHITE


def _add_critical_equipment_slide(prs: Presentation, report_data: dict):
    """Slide 5: Top 10 Critical Equipment table."""
    slides_data = _safe_get(report_data, "slides_data", [])
    equip_list: list = []
    for sd in slides_data:
        if _safe_get(sd, "type") == "table" and "equipo" in _safe_get(sd, "title", "").lower():
            equip_list = _safe_get(sd, "data", [])
            break
    if not equip_list:
        equip_list = _safe_get(report_data, "critical_equipment", []) or []

    # Limit to 10
    equip_list = equip_list[:10]

    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Top 10 Equipos Criticos"

    num_rows = 1 + max(len(equip_list), 1)
    cols = 5
    table_shape = slide.shapes.add_table(num_rows, cols, Inches(0.3), Inches(1.8),
                                         Inches(9.4), Inches(0.35 * num_rows + 0.4))
    table = table_shape.table

    # Header
    for i, h in enumerate(["Tag", "Nombre", "Criticidad", "Area", "Ubicacion SAP"]):
        _set_cell_text(table.cell(0, i), h, font_size=10, bold=True)
    _style_header_row(table, cols, font_size=10)

    if not equip_list:
        _set_cell_text(table.cell(1, 0), "Sin datos disponibles", font_size=10,
                       alignment=PP_ALIGN.LEFT)
        return

    for row_idx, eq in enumerate(equip_list, start=1):
        _set_cell_text(table.cell(row_idx, 0), _safe_get(eq, "tag", "N/A"),
                       font_size=9, alignment=PP_ALIGN.LEFT)
        _set_cell_text(table.cell(row_idx, 1), _safe_get(eq, "name", "N/A"),
                       font_size=9, alignment=PP_ALIGN.LEFT)
        crit = _safe_get(eq, "criticality", "N/A")
        _set_cell_text(table.cell(row_idx, 2), str(crit), font_size=9, bold=True)
        _set_cell_text(table.cell(row_idx, 3), _safe_get(eq, "area", "N/A"),
                       font_size=9, alignment=PP_ALIGN.LEFT)
        _set_cell_text(table.cell(row_idx, 4), _safe_get(eq, "sap_func_loc", "N/A"),
                       font_size=9, alignment=PP_ALIGN.LEFT)

        # Criticality colour
        crit_color = _criticality_color(str(crit))
        if crit_color:
            _set_cell_bg(table.cell(row_idx, 2), crit_color)
            for paragraph in table.cell(row_idx, 2).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = WHITE

    _apply_alternating_rows(table, num_rows, cols)
    # Re-apply criticality colours
    for row_idx, eq in enumerate(equip_list, start=1):
        crit = _safe_get(eq, "criticality", "")
        crit_color = _criticality_color(str(crit))
        if crit_color:
            _set_cell_bg(table.cell(row_idx, 2), crit_color)
            for paragraph in table.cell(row_idx, 2).text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = WHITE


def _add_text_slide(prs: Presentation, title: str, body_text: str):
    """Generic text slide (used for summary and next actions)."""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title

    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(9.0), Inches(5.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(body_text) if body_text else "Sin datos disponibles."
    p.font.size = Pt(14)
    p.font.color.rgb = DARK_BLUE
    p.alignment = PP_ALIGN.LEFT


def _add_closing_slide(prs: Presentation):
    """Slide 8: Closing / Thank you."""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Gracias"
    if slide.placeholders[1]:
        slide.placeholders[1].text = "Generado por AMS \u2014 Agentic Solutions"


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

def generate_pptx(report_data: dict) -> str:
    """Generate a PPTX file from executive report data.

    Parameters
    ----------
    report_data : dict
        Output of ``generate_executive_report`` containing summary_text,
        slides_data, kpis, backlog_stats, budget_variance, critical_equipment,
        period, plant_id, and generated_at.

    Returns
    -------
    str
        Absolute filepath to the generated .pptx file.
    """
    log.info("Generating executive PPTX report")
    report_data = report_data or {}

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # -- Build slides -------------------------------------------------------
    _add_title_slide(prs, report_data)
    _add_kpi_slide(prs, report_data)
    _add_backlog_slide(prs, report_data)
    _add_budget_slide(prs, report_data)
    _add_critical_equipment_slide(prs, report_data)

    # Text slides: summary and next actions
    slides_data = _safe_get(report_data, "slides_data", [])
    summary_text = _safe_get(report_data, "summary_text", "")
    next_actions_text = ""
    for sd in slides_data:
        if _safe_get(sd, "type") == "text":
            title_lower = _safe_get(sd, "title", "").lower()
            if "resumen" in title_lower or "summary" in title_lower:
                summary_text = summary_text or _safe_get(sd, "data", "")
            elif "accion" in title_lower or "action" in title_lower or "proxima" in title_lower:
                next_actions_text = _safe_get(sd, "data", "")

    _add_text_slide(prs, "Resumen Ejecutivo", summary_text)
    _add_text_slide(prs, "Proximas Acciones", next_actions_text)
    _add_closing_slide(prs)

    # -- Save file ----------------------------------------------------------
    plant_id = _safe_get(report_data, "plant_id", "unknown")
    date_str = datetime.now().strftime("%Y%m%d")
    filename = f"executive_report_{plant_id}_{date_str}.pptx"
    tmp_dir = tempfile.mkdtemp()
    filepath = os.path.join(tmp_dir, filename)
    prs.save(filepath)

    log.info("PPTX report saved to %s", filepath)
    return filepath


def generate_pptx_response(report_data: dict):
    """Generate PPTX and return a FastAPI ``FileResponse``.

    Convenience wrapper for use in API routers.
    """
    from fastapi.responses import FileResponse

    filepath = generate_pptx(report_data)
    filename = os.path.basename(filepath)
    return FileResponse(
        filepath,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename=filename,
    )
