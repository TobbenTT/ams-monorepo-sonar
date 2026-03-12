"""OR System document generation tools: Word, Excel, PowerPoint, and deliverable listing.

Used by all OR agents that need to produce professional deliverables.
Requires python-docx, python-pptx, and openpyxl (all in requirements.txt).
"""

import json
import os
import uuid
from datetime import datetime
from pathlib import Path

from api.ai_core.tool_wrappers.registry import tool

OR_DELIVERABLES_DIR = os.getenv("OR_DELIVERABLES_DIR", "/app/or_deliverables")


def _ensure_dir() -> Path:
    p = Path(OR_DELIVERABLES_DIR)
    p.mkdir(parents=True, exist_ok=True)
    return p


# ── Word Document ────────────────────────────────────────────────────────────

@tool(
    "generate_word_document",
    "Generate a professional Word (.docx) document with title, sections, and content.",
    {
        "type": "object",
        "properties": {
            "input_json": {
                "type": "string",
                "description": (
                    "JSON with: title (str), sections (list of {heading, content}), "
                    "author (str, optional), project_id (str, optional)"
                ),
            }
        },
        "required": ["input_json"],
    },
)
def generate_word_document(input_json: str) -> str:
    try:
        data = json.loads(input_json)
    except json.JSONDecodeError as e:
        return json.dumps({"error": f"Invalid JSON: {e}"})

    title = data.get("title", "Document")
    sections = data.get("sections", [])
    author = data.get("author", "CORTEX")
    project_id = data.get("project_id", "")

    try:
        from docx import Document
        from docx.shared import Pt, RGBColor
        from docx.enum.text import WD_ALIGN_PARAGRAPH

        doc = Document()

        # Title
        title_para = doc.add_heading(title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER

        # Metadata
        meta = doc.add_paragraph()
        meta.add_run(f"Author: {author}  |  Generated: {datetime.utcnow().strftime('%Y-%m-%d %H:%M UTC')}")
        meta.alignment = WD_ALIGN_PARAGRAPH.CENTER
        doc.add_paragraph()

        # Sections
        for sec in sections:
            heading = sec.get("heading", "Section")
            content = sec.get("content", "")
            doc.add_heading(heading, level=1)
            if isinstance(content, list):
                for item in content:
                    doc.add_paragraph(str(item), style="List Bullet")
            else:
                doc.add_paragraph(str(content))
            doc.add_paragraph()

        # Save
        out_dir = _ensure_dir()
        filename = f"{uuid.uuid4().hex[:8]}_{title[:30].replace(' ', '_')}.docx"
        filepath = out_dir / filename
        doc.save(str(filepath))

        return json.dumps({
            "status": "ok",
            "filename": filename,
            "path": str(filepath),
            "size_bytes": filepath.stat().st_size,
            "project_id": project_id,
        })

    except ImportError:
        return json.dumps({
            "error": "python-docx not installed. Run: pip install python-docx",
            "title": title,
            "sections_count": len(sections),
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


# ── Excel Workbook ───────────────────────────────────────────────────────────

@tool(
    "generate_excel_workbook",
    "Generate an Excel (.xlsx) workbook with one or more sheets of tabular data.",
    {
        "type": "object",
        "properties": {
            "input_json": {
                "type": "string",
                "description": (
                    "JSON with: title (str), sheets (list of {name, headers, rows}), "
                    "project_id (str, optional)"
                ),
            }
        },
        "required": ["input_json"],
    },
)
def generate_excel_workbook(input_json: str) -> str:
    try:
        data = json.loads(input_json)
    except json.JSONDecodeError as e:
        return json.dumps({"error": f"Invalid JSON: {e}"})

    title = data.get("title", "Workbook")
    sheets = data.get("sheets", [])
    project_id = data.get("project_id", "")

    try:
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment
        from openpyxl.utils import get_column_letter

        wb = openpyxl.Workbook()
        wb.remove(wb.active)  # Remove default sheet

        header_fill = PatternFill("solid", fgColor="1F4E79")
        header_font = Font(color="FFFFFF", bold=True)

        for sheet_data in sheets:
            ws = wb.create_sheet(title=sheet_data.get("name", "Sheet")[:31])
            headers = sheet_data.get("headers", [])
            rows = sheet_data.get("rows", [])

            # Headers
            for col_idx, header in enumerate(headers, 1):
                cell = ws.cell(row=1, column=col_idx, value=header)
                cell.font = header_font
                cell.fill = header_fill
                cell.alignment = Alignment(horizontal="center")

            # Data rows
            for row_idx, row in enumerate(rows, 2):
                if isinstance(row, dict):
                    row_values = [row.get(h, "") for h in headers]
                elif isinstance(row, list):
                    row_values = row
                else:
                    row_values = [str(row)]
                for col_idx, val in enumerate(row_values, 1):
                    ws.cell(row=row_idx, column=col_idx, value=val)

            # Auto-width
            for col_idx in range(1, len(headers) + 1):
                ws.column_dimensions[get_column_letter(col_idx)].width = 20

        if not wb.sheetnames:
            wb.create_sheet("Data")

        out_dir = _ensure_dir()
        filename = f"{uuid.uuid4().hex[:8]}_{title[:30].replace(' ', '_')}.xlsx"
        filepath = out_dir / filename
        wb.save(str(filepath))

        return json.dumps({
            "status": "ok",
            "filename": filename,
            "path": str(filepath),
            "size_bytes": filepath.stat().st_size,
            "sheets": [s.get("name", "Sheet") for s in sheets],
            "project_id": project_id,
        })

    except ImportError:
        return json.dumps({
            "error": "openpyxl not installed. Run: pip install openpyxl",
            "title": title,
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


# ── PowerPoint Presentation ──────────────────────────────────────────────────

@tool(
    "generate_presentation",
    "Generate a PowerPoint (.pptx) presentation with slides containing titles and bullet points.",
    {
        "type": "object",
        "properties": {
            "input_json": {
                "type": "string",
                "description": (
                    "JSON with: title (str), slides (list of {title, bullets or content}), "
                    "project_id (str, optional)"
                ),
            }
        },
        "required": ["input_json"],
    },
)
def generate_presentation(input_json: str) -> str:
    try:
        data = json.loads(input_json)
    except json.JSONDecodeError as e:
        return json.dumps({"error": f"Invalid JSON: {e}"})

    title = data.get("title", "Presentation")
    slides_data = data.get("slides", [])
    project_id = data.get("project_id", "")

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        prs = Presentation()
        blank_layout = prs.slide_layouts[1]  # Title and Content

        # Title slide
        title_layout = prs.slide_layouts[0]
        title_slide = prs.slides.add_slide(title_layout)
        title_slide.shapes.title.text = title
        if len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = f"Generated: {datetime.utcnow().strftime('%Y-%m-%d')}"

        # Content slides
        for slide_data in slides_data:
            slide = prs.slides.add_slide(blank_layout)
            slide.shapes.title.text = slide_data.get("title", "Slide")

            content_placeholder = slide.placeholders[1]
            tf = content_placeholder.text_frame
            tf.word_wrap = True

            bullets = slide_data.get("bullets") or slide_data.get("content", [])
            if isinstance(bullets, str):
                bullets = [bullets]
            elif isinstance(bullets, list):
                pass
            else:
                bullets = [str(bullets)]

            for i, bullet in enumerate(bullets):
                if i == 0:
                    tf.text = str(bullet)
                else:
                    p = tf.add_paragraph()
                    p.text = str(bullet)

        out_dir = _ensure_dir()
        filename = f"{uuid.uuid4().hex[:8]}_{title[:30].replace(' ', '_')}.pptx"
        filepath = out_dir / filename
        prs.save(str(filepath))

        return json.dumps({
            "status": "ok",
            "filename": filename,
            "path": str(filepath),
            "size_bytes": filepath.stat().st_size,
            "slides_count": len(slides_data) + 1,
            "project_id": project_id,
        })

    except ImportError:
        return json.dumps({
            "error": "python-pptx not installed. Run: pip install python-pptx",
            "title": title,
        })
    except Exception as e:
        return json.dumps({"error": str(e)})


# ── List Deliverables ────────────────────────────────────────────────────────

@tool(
    "list_deliverables",
    "List all generated deliverable files (docx, xlsx, pptx) in the OR deliverables directory.",
    {
        "type": "object",
        "properties": {
            "input_json": {
                "type": "string",
                "description": "JSON with optional: file_type (docx/xlsx/pptx), project_id (str filter)",
            }
        },
        "required": ["input_json"],
    },
)
def list_deliverables(input_json: str) -> str:
    try:
        data = json.loads(input_json) if input_json.strip() else {}
    except json.JSONDecodeError:
        data = {}

    file_type_filter = data.get("file_type", "").lower()

    out_dir = Path(OR_DELIVERABLES_DIR)
    if not out_dir.exists():
        return json.dumps({"files": [], "count": 0, "directory": str(out_dir)})

    extensions = [".docx", ".xlsx", ".pptx"]
    if file_type_filter in ("docx", "xlsx", "pptx"):
        extensions = [f".{file_type_filter}"]

    files = []
    for ext in extensions:
        for f in sorted(out_dir.glob(f"*{ext}"), key=lambda x: x.stat().st_mtime, reverse=True):
            files.append({
                "filename": f.name,
                "path": str(f),
                "type": f.suffix.lstrip("."),
                "size_bytes": f.stat().st_size,
                "modified": datetime.fromtimestamp(f.stat().st_mtime).strftime("%Y-%m-%d %H:%M"),
            })

    return json.dumps({"files": files, "count": len(files), "directory": str(out_dir)})
